from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
ASSETS_DIR = BASE_DIR / "assets"
OUTPUT_PATH = BASE_DIR / "ASEM-Talint-10-minute-presentation.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG = RGBColor(246, 241, 232)
SURFACE = RGBColor(255, 252, 247)
SURFACE_ALT = RGBColor(251, 247, 241)
INK = RGBColor(35, 32, 28)
MUTED = RGBColor(107, 101, 94)
ACCENT = RGBColor(122, 102, 75)
WARM = RGBColor(214, 69, 45)
WARM_SOFT = RGBColor(242, 201, 76)
BORDER = RGBColor(223, 214, 201)
DARK = RGBColor(45, 42, 38)
GREEN = RGBColor(76, 125, 82)


def set_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_round_box(slide, left, top, width, height, fill_color, line_color=BORDER, radius_shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(radius_shape, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = Pt(1)
    return shape


def add_text(slide, left, top, width, height, text, *, font_size=18, color=INK, bold=False, name="Aptos", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = name
    font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, *, font_size=18, color=INK, bullet_color=ACCENT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()
    for index, item in enumerate(items):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.name = "Aptos"
        p.font.color.rgb = color
        p.bullet = True
        p.space_after = Pt(6)
    return box


def add_title_block(slide, eyebrow, title, subtitle=None):
    add_text(slide, Inches(0.7), Inches(0.42), Inches(5.3), Inches(0.35), eyebrow, font_size=12, color=ACCENT, bold=True, name="Aptos")
    add_text(slide, Inches(0.7), Inches(0.78), Inches(5.8), Inches(1.4), title, font_size=28, color=INK, bold=True, name="Aptos Display")
    if subtitle:
        add_text(slide, Inches(0.7), Inches(1.72), Inches(5.8), Inches(0.7), subtitle, font_size=14, color=MUTED, name="Aptos")


def add_footer(slide, text):
    add_text(slide, Inches(0.7), Inches(7.05), Inches(12), Inches(0.25), text, font_size=9, color=MUTED, name="Aptos")


def add_picture_fit(slide, path: Path, left, top, width, height):
    if not path.exists():
        return None
    image = Image.open(path)
    image_ratio = image.width / image.height
    box_ratio = width / height
    if image_ratio > box_ratio:
        pic_width = width
        pic_height = width / image_ratio
        pic_left = left
        pic_top = top + (height - pic_height) / 2
    else:
        pic_height = height
        pic_width = height * image_ratio
        pic_top = top
        pic_left = left + (width - pic_width) / 2
    return slide.shapes.add_picture(str(path), pic_left, pic_top, width=pic_width, height=pic_height)


def add_card_with_title(slide, left, top, width, height, title, lines, *, fill_color=SURFACE):
    add_round_box(slide, left, top, width, height, fill_color)
    add_text(slide, left + Inches(0.18), top + Inches(0.14), width - Inches(0.36), Inches(0.28), title, font_size=11, color=ACCENT, bold=True, name="Aptos")
    add_bullets(slide, left + Inches(0.18), top + Inches(0.48), width - Inches(0.32), height - Inches(0.56), lines, font_size=14)


def build_slide_1(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_round_box(slide, Inches(0.45), Inches(0.35), Inches(12.4), Inches(6.8), SURFACE)
    add_text(slide, Inches(0.85), Inches(0.72), Inches(3.6), Inches(0.32), "UM Hackathon 2026 | Domain 2", font_size=12, color=ACCENT, bold=True)
    add_text(slide, Inches(0.85), Inches(1.05), Inches(5.2), Inches(1.35), "ASEM Talint", font_size=30, color=INK, bold=True, name="Aptos Display")
    add_text(slide, Inches(0.85), Inches(2.02), Inches(5.2), Inches(1.2), "Dashboard-first semiconductor talent decision intelligence", font_size=20, color=DARK, bold=True)
    add_bullets(
        slide,
        Inches(0.85),
        Inches(2.85),
        Inches(5.0),
        Inches(1.9),
        [
            "Deterministic readiness scoring plus Z.AI GLM structured reasoning",
            "Built for training-path fit, evidence review, and market-context visibility",
            "Team Novum: Ong Khai Sern (team leader), Tan Eng Feng (team member)",
        ],
        font_size=16,
    )
    add_round_box(slide, Inches(6.1), Inches(0.78), Inches(6.1), Inches(5.5), SURFACE_ALT)
    add_picture_fit(slide, ASSETS_DIR / "dashboard-home.png", Inches(6.2), Inches(0.9), Inches(5.9), Inches(5.25))
    add_round_box(slide, Inches(0.85), Inches(5.25), Inches(5.0), Inches(0.8), SURFACE_ALT)
    add_text(
        slide,
        Inches(1.05),
        Inches(5.48),
        Inches(4.6),
        Inches(0.32),
        "One place to read the candidate clearly, with live market context and auditable AI explanation.",
        font_size=14,
        color=DARK,
        bold=True,
    )
    add_footer(slide, "ASEM Talint preliminary deck | Team Novum | Local prototype screenshots from current build")


def build_slide_2(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title_block(slide, "Slide 2", "Problem: placement decisions are still fragmented", "Show domain fit, not just technical ambition.")
    add_card_with_title(
        slide,
        Inches(0.7), Inches(1.55), Inches(4.0), Inches(4.7),
        "What breaks today",
        [
            "Candidate-to-track decisions are still reviewed manually and inconsistently.",
            "Missing inputs, access constraints, and labor-market context are often separated.",
            "Pathway decisions should connect training to real economic outcomes, not only admissions workflow.",
        ],
    )
    add_card_with_title(
        slide,
        Inches(4.9), Inches(1.55), Inches(3.45), Inches(4.7),
        "Why it matters",
        [
            "Semiconductor pathways require scarce seats, scarce OJT capacity, and better matching discipline.",
            "Reviewers need auditable evidence, not a generic chat response.",
            "Malaysia's industrial push makes local-value pathway allocation more strategic.",
        ],
    )
    add_card_with_title(
        slide,
        Inches(8.55), Inches(1.55), Inches(4.05), Inches(4.7),
        "Our response",
        [
            "Standardize candidate evidence.",
            "Combine deterministic scoring with official data slices.",
            "Use Z.AI GLM for structured explanation on the judged path.",
        ],
        fill_color=RGBColor(255, 248, 239),
    )
    add_footer(slide, "Citations: UMHackathon Domain 2 brief; UMHackathon 2026 Judging Criteria")


def build_slide_3(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title_block(slide, "Slide 3", "Economic theory: what friction we reduce", "Anchor the product in economic logic, not only features.")
    titles = [
        ("Human capital allocation", "Training seats and OJT slots are scarce investments."),
        ("Signaling and information asymmetry", "Resumes and credentials are noisy until structured."),
        ("Search and matching frictions", "Candidate, track, and employer fit must be aligned."),
        ("Spatial and wage frictions", "Access and economic payoff change pathway quality."),
    ]
    positions = [(0.7, 1.6), (6.7, 1.6), (0.7, 3.55), (6.7, 3.55)]
    for (title, text), (left, top) in zip(titles, positions, strict=True):
        add_round_box(slide, Inches(left), Inches(top), Inches(5.9), Inches(1.55), SURFACE)
        add_text(slide, Inches(left + 0.22), Inches(top + 0.2), Inches(5.4), Inches(0.35), title, font_size=18, color=DARK, bold=True)
        add_text(slide, Inches(left + 0.22), Inches(top + 0.66), Inches(5.35), Inches(0.52), text, font_size=14, color=MUTED)
    add_round_box(slide, Inches(0.7), Inches(5.58), Inches(11.95), Inches(0.72), RGBColor(255, 248, 239))
    add_text(slide, Inches(0.95), Inches(5.84), Inches(11.4), Inches(0.22), "Feature mapping: Candidate Lab -> signals | Market Studio/OJT -> matching | GIS accessibility -> spatial friction | Wage mobility -> economic payoff | ERP Bridge -> coordination cost", font_size=14, color=DARK, bold=True)
    add_footer(slide, "Economic framing source: submission/preliminary/economic-theory-anchor.md")


def build_slide_4(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title_block(slide, "Slide 4", "Why this matters in Malaysia now", "Tie the product to current industrial capability-building logic.")
    add_round_box(slide, Inches(0.7), Inches(1.58), Inches(5.2), Inches(4.65), SURFACE)
    add_bullets(
        slide,
        Inches(0.95), Inches(1.95), Inches(4.7), Inches(3.8),
        [
            '"Made by Malaysia" raises the bar for local value creation.',
            "Talent pathways need stronger links to vendors, SMEs, and industry demand.",
            "Training decisions should be explainable, evidence-backed, and fast.",
            "A pathway is stronger when access, wage direction, and employer demand all support it.",
        ],
        font_size=17,
    )
    add_round_box(slide, Inches(6.15), Inches(1.58), Inches(6.45), Inches(4.65), RGBColor(255, 248, 239), line_color=RGBColor(230, 198, 167))
    add_text(slide, Inches(6.45), Inches(1.98), Inches(5.9), Inches(0.35), "Malaysia lens", font_size=12, color=ACCENT, bold=True)
    add_text(slide, Inches(6.45), Inches(2.28), Inches(5.75), Inches(0.9), "Local capability, not just placement", font_size=24, color=DARK, bold=True, name="Aptos Display")
    add_text(slide, Inches(6.45), Inches(3.12), Inches(5.6), Inches(1.8), "The national push is toward stronger local vendors, deeper SME participation, and tighter links between investment, training, and research. ASEM Talint keeps that lens visible while pathway decisions are being made.", font_size=16, color=MUTED)
    add_text(slide, Inches(6.45), Inches(5.2), Inches(4.8), Inches(0.3), "NST source, 5 Feb 2026", font_size=11, color=WARM, bold=True)
    add_footer(slide, "Citation: New Straits Times, 5 Feb 2026, 'Shift to Made by Malaysia strategy'")


def build_slide_5(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title_block(slide, "Slide 5", "Users and decision workflow", "Make the target users and the review loop concrete.")
    add_card_with_title(slide, Inches(0.7), Inches(1.65), Inches(3.65), Inches(2.05), "Primary user", ["Training reviewer or program lead", "Needs fit, gaps, readiness, and next action in one place"], fill_color=SURFACE)
    add_card_with_title(slide, Inches(0.7), Inches(3.95), Inches(3.65), Inches(2.05), "Supporting user", ["Employer-partnership or OJT reviewer", "Needs match quality, access, and wage context"], fill_color=SURFACE)
    stages = [
        "1. Ingest candidate and track evidence",
        "2. Parse resume and standardize signals",
        "3. Compute deterministic fit and market context",
        "4. Use Z.AI GLM to generate structured explanation",
    ]
    for idx, stage in enumerate(stages):
        left = 4.8 + idx * 1.92
        add_round_box(slide, Inches(left), Inches(2.2), Inches(1.7), Inches(2.4), SURFACE_ALT)
        add_text(slide, Inches(left + 0.16), Inches(2.45), Inches(1.38), Inches(1.7), stage, font_size=14, color=DARK, bold=True)
        if idx < len(stages) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(left + 1.62), Inches(3.0), Inches(0.28), Inches(0.62))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = WARM_SOFT
            arrow.line.color.rgb = WARM_SOFT
    add_footer(slide, "Workflow source: PRD, pitch deck outline, and current live prototype flow")


def build_slide_6(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title_block(slide, "Slide 6", "MVP scope and product surfaces", "Demonstrate feature prioritization and show real snapshots from the current build.")
    add_card_with_title(slide, Inches(0.7), Inches(1.6), Inches(3.65), Inches(2.05), "In scope now", ["One candidate-to-track fit workflow", "Resume parsing for PDF and DOCX with OCR fallback", "Wages, employer demand, accessibility, OJT, and wage-mobility panels"], fill_color=SURFACE)
    add_card_with_title(slide, Inches(0.7), Inches(3.95), Inches(3.65), Inches(2.05), "Still out of scope", ["Production authentication and reviewer case management", "Full persistence and operational analytics", "Public deployment claims beyond local validation"], fill_color=SURFACE)
    placements = [
        (ASSETS_DIR / "candidate-lab.png", 4.7, 1.68),
        (ASSETS_DIR / "market-studio.png", 8.72, 1.68),
        (ASSETS_DIR / "pathway-planner.png", 4.7, 4.1),
        (ASSETS_DIR / "erp-bridge.png", 8.72, 4.1),
    ]
    for image_path, left, top in placements:
        add_round_box(slide, Inches(left), Inches(top), Inches(3.45), Inches(2.0), SURFACE_ALT)
        add_picture_fit(slide, image_path, Inches(left + 0.05), Inches(top + 0.05), Inches(3.35), Inches(1.9))
    add_footer(slide, "Snapshots: Candidate Lab, Market Studio, Pathway Planner, ERP Bridge from current local app")


def build_slide_7(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title_block(slide, "Slide 7", "Why Z.AI GLM is central", "Address judging eligibility and technical intent directly.")
    add_round_box(slide, Inches(0.7), Inches(1.65), Inches(12.0), Inches(4.9), SURFACE)
    components = [
        ("Deterministic context", "Scoring, thresholds, missing-input handling, and wage/access logic remain in code.", 0.95),
        ("Prompt builder", "Bounded JSON evidence is passed forward so the model explains a structured decision problem.", 3.45),
        ("Z.AI GLM provider", "This is the judged reasoning path in the shipped product flow.", 5.95),
        ("Schema validation", "Model output must parse into structured response JSON before the app uses it.", 8.45),
        ("Dashboard/API consumer", "Reviewers see recommendation, tradeoffs, missing inputs, and pathway steps.", 10.95),
    ]
    for title, body, left in components:
        add_round_box(slide, Inches(left), Inches(2.6), Inches(1.9), Inches(2.4), SURFACE_ALT)
        add_text(slide, Inches(left + 0.12), Inches(2.82), Inches(1.66), Inches(0.45), title, font_size=16, color=DARK, bold=True)
        add_text(slide, Inches(left + 0.12), Inches(3.26), Inches(1.66), Inches(1.45), body, font_size=12, color=MUTED)
    for x in [2.78, 5.28, 7.78, 10.28]:
        arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(3.35), Inches(0.38), Inches(0.62))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = WARM
        arrow.line.color.rgb = WARM
    add_round_box(slide, Inches(0.95), Inches(5.35), Inches(11.5), Inches(0.62), RGBColor(255, 248, 239), line_color=RGBColor(230, 198, 167))
    add_text(slide, Inches(1.12), Inches(5.56), Inches(11.1), Inches(0.2), "No silent model replacement: the optional ILMU route is compatibility infrastructure, not the judged runtime path.", font_size=14, color=DARK, bold=True)
    add_footer(slide, "Architecture rule: keep Z.AI GLM central in implementation, narrative, and demo")


def build_slide_8(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title_block(slide, "Slide 8", "System architecture and feasibility", "Explain why the prototype is technically coherent and operationally plausible.")
    layers = [
        ("FastAPI interface", "Dashboard, decision API, resume API, market routes"),
        ("Decision engine", "Deterministic context + provider boundary"),
        ("Data and services", "Resume parsing, wages, hotspots, OJT, wage mobility"),
        ("Output contracts", "Structured JSON, ERP package, observable errors"),
    ]
    for idx, (title, body) in enumerate(layers):
        top = 1.7 + idx * 1.2
        add_round_box(slide, Inches(0.9), Inches(top), Inches(4.2), Inches(0.88), SURFACE)
        add_text(slide, Inches(1.12), Inches(top + 0.14), Inches(1.7), Inches(0.26), title, font_size=16, color=DARK, bold=True)
        add_text(slide, Inches(2.85), Inches(top + 0.14), Inches(1.95), Inches(0.4), body, font_size=12, color=MUTED)
    add_round_box(slide, Inches(5.55), Inches(1.7), Inches(6.95), Inches(4.95), SURFACE_ALT)
    add_bullets(
        slide,
        Inches(5.9), Inches(2.05), Inches(6.3), Inches(4.2),
        [
            "Typed request and response models keep the interfaces explicit.",
            "Prompt budgets, token budgets, and upload limits are enforced in code.",
            "Official-data slices are normalized locally for preliminary feasibility.",
            "ERP Bridge packages the same decision output into operational sync-ready JSON.",
            "The architecture reduces screening, matching, spatial, and coordination frictions in one workflow.",
        ],
        font_size=16,
    )
    add_footer(slide, "Feasibility source: SAD, README, and live prototype implementation")


def build_slide_9(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title_block(slide, "Slide 9", "Live demo walkthrough", "Show the prototype rather than only talk about it.")
    add_bullets(
        slide,
        Inches(0.8), Inches(1.7), Inches(5.0), Inches(4.9),
        [
            "Show runtime-ready status and the sample candidate payload.",
            "Parse and apply a resume so judges see structured evidence ingestion.",
            "Refresh market panels for wages, employer demand, accessibility, and OJT.",
            "Pause on the cleaned GIS accessibility figure with numbered hotspots and ranked list.",
            "Run the live Z.AI route and inspect recommendation, tradeoffs, and missing inputs.",
            "Open the raw structured response briefly to prove contract discipline.",
        ],
        font_size=17,
    )
    add_round_box(slide, Inches(6.05), Inches(1.65), Inches(6.2), Inches(4.85), SURFACE_ALT)
    add_picture_fit(slide, ASSETS_DIR / "dashboard-gis-accessibility-panel.png", Inches(6.15), Inches(1.78), Inches(6.0), Inches(4.55))
    add_footer(slide, "Fallback still image in deck: cleaned GIS accessibility panel from current dashboard build")


def build_slide_10(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title_block(slide, "Slide 10", "Engineering quality and risk control", "Answer code quality and testing feasibility criteria.")
    add_card_with_title(slide, Inches(0.7), Inches(1.6), Inches(3.7), Inches(2.15), "Quality controls", ["Typed models and explicit provider boundaries", "Provider output validated before use", "Malformed or empty model output fails visibly"], fill_color=SURFACE)
    add_card_with_title(slide, Inches(4.8), Inches(1.6), Inches(3.7), Inches(2.15), "Test coverage", ["Happy path and invalid payload routes", "OCR fallback and parser warnings", "Health/readiness behavior and judged-path guardrails"], fill_color=SURFACE)
    add_card_with_title(slide, Inches(8.9), Inches(1.6), Inches(3.7), Inches(2.15), "Observability", ["Provider usage metrics surfaced", "Warnings kept explicit", "No silent degradation of judged model path"], fill_color=SURFACE)
    add_round_box(slide, Inches(0.7), Inches(4.2), Inches(11.9), Inches(2.1), SURFACE_ALT)
    add_text(slide, Inches(0.95), Inches(4.45), Inches(2.4), Inches(0.25), "Risk excerpt", font_size=12, color=ACCENT, bold=True)
    risk_rows = [
        ("Live provider readiness", "Medium", "Validate `/health`, fail visibly, no silent substitution"),
        ("Data refresh discipline", "Medium", "Use bounded official-data slices and document refresh path"),
        ("Submission packaging quality", "Low-Med", "Keep PRD, SAD, QATD, repo, deck, and video aligned"),
    ]
    for idx, row in enumerate(risk_rows):
        top = 4.82 + idx * 0.42
        add_text(slide, Inches(1.0), Inches(top), Inches(2.2), Inches(0.22), row[0], font_size=13, color=DARK, bold=True)
        add_text(slide, Inches(3.42), Inches(top), Inches(1.0), Inches(0.22), row[1], font_size=13, color=WARM, bold=True)
        add_text(slide, Inches(4.55), Inches(top), Inches(7.4), Inches(0.28), row[2], font_size=13, color=MUTED)
    add_footer(slide, "Validation evidence: pytest route tests, provider parsing, OCR fallback, and documented QATD risk controls")


def build_slide_11(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_round_box(slide, Inches(0.45), Inches(0.45), Inches(12.4), Inches(6.55), SURFACE)
    add_text(slide, Inches(0.85), Inches(0.82), Inches(5.4), Inches(0.32), "Closing", font_size=12, color=ACCENT, bold=True)
    add_text(slide, Inches(0.85), Inches(1.12), Inches(5.6), Inches(1.05), "Feasible prototype, honest scope, submission-ready structure", font_size=26, color=INK, bold=True, name="Aptos Display")
    add_bullets(
        slide,
        Inches(0.85), Inches(2.15), Inches(5.4), Inches(2.9),
        [
            "ASEM Talint fits Domain 2 because it supports economically meaningful pathway decisions.",
            "Z.AI GLM remains central in the product, architecture, and demo narrative.",
            "PRD, SAD, QATD, repository, deck, and video runbook are structured for preliminary submission.",
            "Repository: github.com/khaisernong/ASEM-Talint",
        ],
        font_size=17,
    )
    collage = [
        (ASSETS_DIR / "dashboard-home.png", 6.45, 0.95, 2.7, 2.0),
        (ASSETS_DIR / "candidate-lab.png", 9.3, 0.95, 2.7, 2.0),
        (ASSETS_DIR / "dashboard-gis-accessibility-panel.png", 6.45, 3.15, 2.7, 2.0),
        (ASSETS_DIR / "erp-bridge.png", 9.3, 3.15, 2.7, 2.0),
    ]
    for image_path, left, top, width, height in collage:
        add_round_box(slide, Inches(left), Inches(top), Inches(width), Inches(height), SURFACE_ALT)
        add_picture_fit(slide, image_path, Inches(left + 0.04), Inches(top + 0.04), Inches(width - 0.08), Inches(height - 0.08))
    add_text(slide, Inches(0.85), Inches(5.62), Inches(5.6), Inches(0.5), "ASEM Talint turns candidate placement from a manual review task into an auditable, market-aware, Z.AI-centered decision workflow.", font_size=16, color=DARK, bold=True)
    add_footer(slide, "Team Novum | Ong Khai Sern | Tan Eng Feng | Preliminary round presentation deck")


def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    build_slide_1(prs)
    build_slide_2(prs)
    build_slide_3(prs)
    build_slide_4(prs)
    build_slide_5(prs)
    build_slide_6(prs)
    build_slide_7(prs)
    build_slide_8(prs)
    build_slide_9(prs)
    build_slide_10(prs)
    build_slide_11(prs)
    prs.save(str(OUTPUT_PATH))
    return OUTPUT_PATH


if __name__ == "__main__":
    output = build_presentation()
    print(f"Generated {output}")